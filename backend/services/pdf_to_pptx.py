"""
services/pdf_to_pptx.py
Converts a PDF to a PowerPoint file (.pptx) by rendering each page as
a full-bleed slide image.

Approach: page-as-image (reliable, looks identical to source).
True text/shape reconstruction is left as a future enhancement.
Uses PyMuPDF for rendering and python-pptx for slide assembly.
"""
import asyncio
import logging
from io import BytesIO
from pathlib import Path

logger = logging.getLogger(__name__)

# Slide dimensions matching a standard 16:9 widescreen (in EMUs)
# 1 inch = 914400 EMUs; 10 × 7.5 inches is the classic 4:3 default.
# We detect the page's own aspect ratio and set slide dimensions to match.
_EMU_PER_PT = 12700   # 1 point = 12700 EMU


def _do_convert(pdf_path: Path, pptx_path: Path, dpi: int) -> int:
    import fitz                         # PyMuPDF
    from pptx import Presentation
    from pptx.util import Emu

    doc = fitz.open(str(pdf_path))
    prs = Presentation()

    # Remove the default blank slide layout (we add our own)
    prs.slide_width  = None  # will be set from first page
    prs.slide_height = None

    mat = fitz.Matrix(dpi / 72, dpi / 72)

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        rect = page.rect                # in points

        # Set presentation dimensions from first page
        if page_num == 0:
            prs.slide_width  = Emu(int(rect.width  * _EMU_PER_PT))
            prs.slide_height = Emu(int(rect.height * _EMU_PER_PT))

        # Render page → PNG bytes in memory
        pix      = page.get_pixmap(matrix=mat, alpha=False)
        img_buf  = BytesIO(pix.tobytes("png"))

        # Add blank slide
        blank_layout = prs.slide_layouts[6]   # "Blank" layout
        slide = prs.slides.add_slide(blank_layout)

        # Add image occupying the full slide area
        slide.shapes.add_picture(
            img_buf,
            left=Emu(0),
            top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

    doc.close()
    prs.save(str(pptx_path))
    return len(doc) if hasattr(doc, '__len__') else page_num + 1


async def pdf_to_pptx(pdf_path: Path, pptx_path: Path, dpi: int = 150) -> int:
    """
    Render each PDF page as a PPTX slide image.

    Args:
        pdf_path:   Source PDF.
        pptx_path:  Destination PPTX.
        dpi:        Render resolution (72 / 150 / 300).

    Returns:
        Number of slides created.
    """
    if dpi not in (72, 150, 300):
        dpi = 150
    loop = asyncio.get_event_loop()
    logger.info(f"PDF→PPTX: {pdf_path.name} at {dpi} DPI")
    count = await loop.run_in_executor(None, _do_convert, pdf_path, pptx_path, dpi)
    logger.info(f"PDF→PPTX complete: {count} slide(s) → {pptx_path.stat().st_size:,} bytes")
    return count
